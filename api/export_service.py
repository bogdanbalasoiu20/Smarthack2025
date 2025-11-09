"""
Export Service for converting presentations to PDF and PowerPoint formats.
"""

import json
import io
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib import colors
from django.http import HttpResponse


class PresentationExportService:
    """Service for exporting presentations to various formats."""

    def __init__(self, presentation):
        self.presentation = presentation
        self.frames = list(presentation.frames.all().order_by('order'))

    def export_to_pptx(self):
        """
        Export presentation to PowerPoint format.

        Returns:
            HttpResponse with PowerPoint file
        """
        prs = PPTXPresentation()
        prs.slide_width = Inches(10)  # 16:9 aspect ratio
        prs.slide_height = Inches(5.625)

        for frame in self.frames:
            # Parse JSON fields
            position = json.loads(frame.position) if isinstance(frame.position, str) else frame.position
            elements = frame.elements.all()

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            bg_color = self._hex_to_rgb(frame.background_color or '#FFFFFF')
            fill.fore_color.rgb = RGBColor(*bg_color)

            # Add elements to slide
            for element in elements:
                element_pos = json.loads(element.position) if isinstance(element.position, str) else element.position
                element_content = json.loads(element.content) if isinstance(element.content, str) else element.content

                if element.element_type == 'TEXT':
                    self._add_text_to_pptx(slide, element_content, element_pos)
                elif element.element_type == 'SHAPE':
                    self._add_shape_to_pptx(slide, element_content, element_pos)

        # Save to BytesIO
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        # Create response
        response = HttpResponse(
            output.read(),
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        filename = f"{self.presentation.title.replace(' ', '_')}.pptx"
        response['Content-Disposition'] = f'attachment; filename="{filename}"'

        return response

    def export_to_pdf(self):
        """
        Export presentation to PDF format.

        Returns:
            HttpResponse with PDF file
        """
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter

        for frame in self.frames:
            # Parse JSON fields
            elements = frame.elements.all()

            # Set background color
            bg_color = self._hex_to_rgb(frame.background_color or '#FFFFFF')
            c.setFillColorRGB(bg_color[0]/255, bg_color[1]/255, bg_color[2]/255)
            c.rect(0, 0, width, height, fill=1, stroke=0)

            # Add elements
            for element in elements:
                element_pos = json.loads(element.position) if isinstance(element.position, str) else element.position
                element_content = json.loads(element.content) if isinstance(element.content, str) else element.content

                if element.element_type == 'TEXT':
                    self._add_text_to_pdf(c, element_content, element_pos, height)
                elif element.element_type == 'SHAPE':
                    self._add_shape_to_pdf(c, element_content, element_pos, height)

            c.showPage()

        c.save()
        buffer.seek(0)

        response = HttpResponse(buffer.read(), content_type='application/pdf')
        filename = f"{self.presentation.title.replace(' ', '_')}.pdf"
        response['Content-Disposition'] = f'attachment; filename="{filename}"'

        return response

    def _add_text_to_pptx(self, slide, content, position):
        """Add text element to PowerPoint slide."""
        # Convert pixel positions to inches (1920x1080 -> 10x5.625 inches)
        left = Inches(position['x'] / 192)
        top = Inches(position['y'] / 192)
        width = Inches(position['width'] / 192)
        height = Inches(position['height'] / 192)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = content.get('text', '')

        # Set font properties
        font = p.font
        font.size = Pt(content.get('fontSize', 24))
        font.name = content.get('fontFamily', 'Arial')

        # Set color
        text_color = self._hex_to_rgb(content.get('color', '#000000'))
        font.color.rgb = RGBColor(*text_color)

        # Set alignment
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT
        }
        p.alignment = align_map.get(content.get('align', 'left'), PP_ALIGN.LEFT)

        # Set bold
        if content.get('fontWeight') == 'bold':
            font.bold = True

    def _add_shape_to_pptx(self, slide, content, position):
        """Add shape element to PowerPoint slide."""
        from pptx.enum.shapes import MSO_SHAPE

        left = Inches(position['x'] / 192)
        top = Inches(position['y'] / 192)
        width = Inches(position['width'] / 192)
        height = Inches(position['height'] / 192)

        shape_type = MSO_SHAPE.OVAL if content.get('shape') == 'circle' else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        # Set fill color
        fill_color = self._hex_to_rgb(content.get('fill', '#818cf8'))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)

        # Set line color
        stroke_color = self._hex_to_rgb(content.get('stroke', '#6366f1'))
        shape.line.color.rgb = RGBColor(*stroke_color)
        shape.line.width = Pt(content.get('strokeWidth', 2))

    def _add_text_to_pdf(self, c, content, position, page_height):
        """Add text element to PDF."""
        # Convert from top-left origin to bottom-left origin
        x = position['x'] * 0.5625  # Scale from 1920 to letter width (612)
        y = page_height - (position['y'] * 0.709)  # Scale and flip Y axis

        # Set font - map to ReportLab built-in fonts
        font_family = content.get('fontFamily', 'Helvetica')
        font_map = {
            'Arial': 'Helvetica',
            'Inter': 'Helvetica',
            'Times New Roman': 'Times-Roman',
            'Courier': 'Courier',
            'Helvetica': 'Helvetica',
        }

        # Get mapped font or default to Helvetica
        pdf_font = font_map.get(font_family, 'Helvetica')

        # Apply bold if needed
        if content.get('fontWeight') == 'bold':
            if pdf_font == 'Helvetica':
                pdf_font = 'Helvetica-Bold'
            elif pdf_font == 'Times-Roman':
                pdf_font = 'Times-Bold'
            elif pdf_font == 'Courier':
                pdf_font = 'Courier-Bold'

        font_size = content.get('fontSize', 24) * 0.7  # Scale font size
        c.setFont(pdf_font, font_size)

        # Set color
        text_color = self._hex_to_rgb(content.get('color', '#000000'))
        c.setFillColorRGB(text_color[0]/255, text_color[1]/255, text_color[2]/255)

        # Draw text
        text = content.get('text', '')

        # Handle alignment (simplified)
        if content.get('align') == 'center':
            x += position['width'] * 0.28  # Approximate center
        elif content.get('align') == 'right':
            x += position['width'] * 0.56  # Approximate right

        # Split long text into lines
        max_width = position['width'] * 0.5625
        words = text.split()
        lines = []
        current_line = []

        for word in words:
            current_line.append(word)
            test_line = ' '.join(current_line)
            if c.stringWidth(test_line, pdf_font, font_size) > max_width:
                if len(current_line) > 1:
                    current_line.pop()
                    lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    lines.append(word)
                    current_line = []

        if current_line:
            lines.append(' '.join(current_line))

        # Draw each line
        for i, line in enumerate(lines):
            c.drawString(x, y - (i * font_size * 1.2), line)

    def _add_shape_to_pdf(self, c, content, position, page_height):
        """Add shape element to PDF."""
        x = position['x'] * 0.5625
        y = page_height - (position['y'] * 0.709) - (position['height'] * 0.709)
        width = position['width'] * 0.5625
        height = position['height'] * 0.709

        # Set fill color
        fill_color = self._hex_to_rgb(content.get('fill', '#818cf8'))
        c.setFillColorRGB(fill_color[0]/255, fill_color[1]/255, fill_color[2]/255)

        # Set stroke color
        stroke_color = self._hex_to_rgb(content.get('stroke', '#6366f1'))
        c.setStrokeColorRGB(stroke_color[0]/255, stroke_color[1]/255, stroke_color[2]/255)
        c.setLineWidth(content.get('strokeWidth', 2))

        # Draw shape
        if content.get('shape') == 'circle':
            c.ellipse(x, y, x + width, y + height, fill=1, stroke=1)
        else:
            c.rect(x, y, width, height, fill=1, stroke=1)

    def _hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
